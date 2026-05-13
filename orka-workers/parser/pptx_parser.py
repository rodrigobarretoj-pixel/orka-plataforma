# parser/pptx_parser.py
"""
Parser estrutural de fichas PowerPoint para ORKA.
Usa python-pptx para extrair textos, tabelas e metadados sem OCR.
Muito superior a OCR para fichamentos criados em PowerPoint.
"""
import io
from pptx import Presentation
from pptx.util import Pt
from typing import Any


def parse_pptx(file_bytes: bytes, filename: str) -> dict[str, Any]:
    """
    Lê um arquivo PPTX e extrai todos os dados estruturados.

    Retorna:
        - slides: lista de slides com textos, tabelas e shapes
        - raw_text: texto bruto completo (para IA)
        - metadata: informações do arquivo
        - tables: tabelas encontradas
        - slide_count: número de slides
    """
    prs = Presentation(io.BytesIO(file_bytes))

    slides_data = []
    all_texts = []
    all_tables = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": slide_num,
            "texts": [],
            "tables": [],
            "shapes": [],
        }

        for shape in slide.shapes:
            shape_data = {
                "type": str(shape.shape_type),
                "name": shape.name,
            }

            # Text frames
            if shape.has_text_frame:
                texts_in_shape = []
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts_in_shape.append(line)

                if texts_in_shape:
                    joined = "\n".join(texts_in_shape)
                    shape_data["text"] = joined
                    slide_info["texts"].append(joined)
                    all_texts.extend(texts_in_shape)

            # Tables
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    if any(row_data):
                        table_data.append(row_data)

                if table_data:
                    shape_data["table"] = table_data
                    slide_info["tables"].append(table_data)
                    all_tables.append({
                        "slide": slide_num,
                        "data": table_data
                    })

            slide_info["shapes"].append(shape_data)

        slides_data.append(slide_info)

    raw_text = "\n\n--- SLIDE ---\n\n".join(
        "\n".join(s["texts"]) for s in slides_data if s["texts"]
    )

    return {
        "filename": filename,
        "slide_count": len(prs.slides),
        "slides": slides_data,
        "tables": all_tables,
        "raw_text": raw_text,
        "metadata": {
            "author": prs.core_properties.author or "",
            "title": prs.core_properties.title or "",
            "subject": prs.core_properties.subject or "",
            "created": str(prs.core_properties.created or ""),
            "modified": str(prs.core_properties.modified or ""),
        }
    }
